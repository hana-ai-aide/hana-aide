"""
K4-01: pptx self-walking walker + importer.
Generates doc.md with ## Slide N sections and <!-- el:<shapeId>/<paraIndex> --> anchors.
Usage:
  python pptx-convert.py --import "<pptx>" "<doc_dir>" [--source "<rel>"]
"""
import os, sys, re, datetime

_EXT_BY_CT = {
    'image/png': '.png', 'image/jpeg': '.jpg', 'image/jpg': '.jpg',
    'image/gif': '.gif', 'image/bmp': '.bmp', 'image/svg+xml': '.svg',
    'image/tiff': '.tif', 'image/x-emf': '.emf', 'image/x-wmf': '.wmf',
    'image/emf': '.emf', 'image/wmf': '.wmf',
}


def _para_text_md(para):
    """Extract text from a pptx paragraph with basic bold/italic markdown decoration."""
    parts = []
    for run in para.runs:
        t = run.text or ''
        if not t:
            continue
        b = run.font.bold is True    # None = inherited, treat as not bold
        i = run.font.italic is True
        if b and i: t = '***' + t + '***'
        elif b:     t = '**' + t + '**'
        elif i:     t = '*' + t + '*'
        parts.append(t)
    return ''.join(parts)


def pptx_import(pptx_path, doc_dir, source_rel=''):
    """K4-01: Walk pptx, generate doc.md with ## Slide N sections and el anchors."""
    from pptx import Presentation

    assets_dir = os.path.join(doc_dir, 'assets')
    os.makedirs(assets_dir, exist_ok=True)

    prs = Presentation(pptx_path)

    img_counter = 0
    img_blob_map = {}   # hash(blob) → 'assets/imageN.ext'
    out_lines = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        out_lines.append('## Slide %d' % slide_num)
        out_lines.append('')

        # Walk shapes in shape_id order for deterministic output
        for shape in sorted(slide.shapes, key=lambda s: s.shape_id):

            # --- Picture shapes ---
            is_picture = False
            try:
                image = shape.image   # AttributeError if not a picture
                is_picture = True
            except AttributeError:
                pass

            if is_picture:
                try:
                    blob = image.blob
                    blob_hash = hash(blob)
                    if blob_hash in img_blob_map:
                        img_ref = img_blob_map[blob_hash]
                    else:
                        ct = image.content_type or ''
                        ext = _EXT_BY_CT.get(ct)
                        if not ext:
                            raw_ext = os.path.splitext(
                                getattr(image, 'ext', '.bin') or '.bin'
                            )[1].lower()
                            ext = raw_ext if raw_ext else '.bin'
                        img_counter += 1
                        fname = 'image%d%s' % (img_counter, ext)
                        with open(os.path.join(assets_dir, fname), 'wb') as f:
                            f.write(blob)
                        img_ref = 'assets/' + fname
                        img_blob_map[blob_hash] = img_ref
                    anchor = '<!-- el:%d/image -->' % shape.shape_id
                    out_lines.append('![](%s)  %s' % (img_ref, anchor))
                    out_lines.append('')
                except Exception as e:
                    print('WARN: image extract failed for shape %d: %s' % (shape.shape_id, e),
                          file=sys.stderr)
                continue

            # --- Text shapes ---
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            for para_idx, para in enumerate(tf.paragraphs):
                text = _para_text_md(para)
                if not text.strip():
                    continue
                anchor = '<!-- el:%d/%d -->' % (shape.shape_id, para_idx)
                out_lines.append('%s  %s' % (text, anchor))
                out_lines.append('')

    # Build body
    body = '\n'.join(out_lines).rstrip() + '\n'
    body = re.sub(r'\n{3,}', '\n\n', body)   # collapse 3+ blank lines

    # Front-matter
    now = datetime.datetime.utcnow().strftime('%Y-%m-%dT%H:%M:%SZ')
    fm = (
        '---\n'
        'type: pptx\n'
        'source: %s\n'
        'convertedAt: %s\n'
        'export: [pptx]\n'
        '---\n\n'
    ) % (source_rel or os.path.basename(pptx_path), now)

    md_path = os.path.join(doc_dir, 'doc.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(fm + body)

    print('MD_FILE: ' + md_path)
    print('ASSETS_DIR: ' + assets_dir)


if __name__ == '__main__':
    if '--import' in sys.argv:
        args = sys.argv[1:]
        ii = args.index('--import')
        rem = args[ii + 1:]
        if len(rem) < 2:
            print('Usage: python pptx-convert.py --import "<pptx>" "<doc_dir>" [--source "<rel>"]')
            sys.exit(1)
        pptx_p = rem[0]
        doc_d  = rem[1]
        src_rel = ''
        if '--source' in rem:
            si = rem.index('--source')
            src_rel = rem[si + 1] if si + 1 < len(rem) else ''
        pptx_import(pptx_p, doc_d, src_rel)
    else:
        print('Usage: python pptx-convert.py --import "<pptx>" "<doc_dir>" [--source "<rel>"]')
        sys.exit(1)
