"""
K4-02: pptx run-level writeback.
Parses doc.md (## Slide N sections + <!-- el:<shapeId>/<paraIndex> --> anchors),
diffs against the base projection, and rewrites only changed a:r runs in source.pptx.
Usage:
  python pptx-writeback.py --pptx <source.pptx> --new <new_doc.md> --base <base_doc.md>
"""
import sys, re, copy, difflib, argparse, json
from lxml import etree

_A   = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
_XML = '{http://www.w3.org/XML/1998/namespace}'

# Anchor: <!-- el:<shapeId>/<paraIndex> -->  (or <shapeId>/image for pictures)
_PPTX_ANCHOR_RE   = re.compile(r'<!--\s*el:(\d+/[^\s>]+)\s*-->', re.I)
_SLIDE_HEADER_RE  = re.compile(r'^##\s+Slide\s+(\d+)\s*$')


# ── MD parsing ─────────────────────────────────────────────────────────────────

def _strip_md_markup(text):
    """Strip common markdown formatting to get plain text for char-level diff."""
    text = _PPTX_ANCHOR_RE.sub('', text)
    text = re.sub(r'^#{1,6}\s+', '', text)
    text = re.sub(r'\*{3}([^*\n]+)\*{3}', r'\1', text)
    text = re.sub(r'\*{2}([^*\n]+)\*{2}', r'\1', text)
    text = re.sub(r'\*([^*\n]+)\*',        r'\1', text)
    return text.strip()


def parse_pptx_md_anchors(md_content):
    """Return dict {composite_key: text} for all anchored text paragraphs.
    composite_key = 'slideN:shapeId/paraIdx'.
    Strips front-matter. Skips ## Slide N headers and image anchors."""
    body = md_content.lstrip()
    if body.startswith('---'):
        end = body.find('\n---\n', 3)
        if end >= 0:
            body = body[end + 5:]

    result = {}
    current_slide = 0
    lines = body.split('\n')

    for i, line in enumerate(lines):
        # Detect ## Slide N
        sh = _SLIDE_HEADER_RE.match(line)
        if sh:
            current_slide = int(sh.group(1))
            continue

        m = _PPTX_ANCHOR_RE.search(line)
        if not m:
            continue

        raw_key = m.group(1)                # e.g. "3/0" or "3/image"
        if raw_key.endswith('/image'):
            continue                        # skip image anchors
        if current_slide == 0:
            continue                        # anchor before any ## Slide header

        composite_key = '%d:%s' % (current_slide, raw_key)  # e.g. "1:3/0"
        line_text = _PPTX_ANCHOR_RE.sub('', line).rstrip()

        if line_text.strip():
            # Inline anchor — text on the same line
            result[composite_key] = line_text.strip()
        else:
            # Standalone anchor — collect non-empty lines immediately above
            block = []
            j = i - 1
            while j >= 0 and lines[j].strip():
                block.insert(0, lines[j])
                j -= 1
            result[composite_key] = '\n'.join(block).strip()

    return result


# ── pptx run-level XML helpers ─────────────────────────────────────────────────

def _pptx_get_run_chars(ap_elem):
    """Flatten <a:r> children into list of (char, a:rPr_or_None)."""
    result = []
    for child in ap_elem:
        if child.tag != _A + 'r':
            continue
        arPr   = child.find(_A + 'rPr')
        t_elem = child.find(_A + 't')
        if t_elem is None or not t_elem.text:
            continue
        for ch in t_elem.text:
            result.append((ch, arPr))
    return result


def _pptx_rpr_key(arPr):
    """Comparable string key for <a:rPr> element."""
    if arPr is None:
        return ''
    return etree.tostring(arPr, encoding='unicode')


def _pptx_make_run(text, arPr):
    """Create an <a:r> lxml element with the given text and arPr."""
    r = etree.Element(_A + 'r')
    if arPr is not None:
        r.append(copy.deepcopy(arPr))
    t = etree.SubElement(r, _A + 't')
    t.text = text
    if text and (text[0] in (' ', '\t') or text[-1] in (' ', '\t')):
        t.set(_XML + 'space', 'preserve')
    return r


def _rewrite_pptx_para(ap_elem, new_plain_text):
    """Rewrite <a:r> runs in an <a:p> element so the combined text equals new_plain_text.
    Unchanged chars keep their original run format; replaced/inserted chars inherit
    the format of the first replaced char (same algorithm as docx-writeback §5 decision 1).
    Returns True if runs were actually changed."""
    old_chars = _pptx_get_run_chars(ap_elem)
    old_text  = ''.join(ch for ch, _ in old_chars)

    if old_text == new_plain_text:
        return False

    # Determine inherit_rPr = rPr of first char being replaced/deleted
    _NOT_SET = object()
    inherit_rPr = _NOT_SET
    matcher  = difflib.SequenceMatcher(None, old_text, new_plain_text, autojunk=False)
    opcodes  = matcher.get_opcodes()
    for tag, i1, i2, j1, j2 in opcodes:
        if tag in ('replace', 'delete') and i1 < len(old_chars):
            inherit_rPr = old_chars[i1][1]
            break
        if tag == 'insert':
            inherit_rPr = old_chars[i1][1] if i1 < len(old_chars) else None
            break
    if inherit_rPr is _NOT_SET:
        inherit_rPr = old_chars[0][1] if old_chars else None

    # Build new (char, rPr) sequence
    new_char_seq = []
    for tag, i1, i2, j1, j2 in opcodes:
        if tag == 'equal':
            new_char_seq.extend(old_chars[i1:i2])
        elif tag in ('replace', 'insert'):
            for ch in new_plain_text[j1:j2]:
                new_char_seq.append((ch, inherit_rPr))
        # delete: skip

    # Group consecutive chars with the same rPr into runs
    runs = []
    if new_char_seq:
        cur_key   = _pptx_rpr_key(new_char_seq[0][1])
        cur_rPr   = new_char_seq[0][1]
        cur_chars = [new_char_seq[0][0]]
        for ch, rPr in new_char_seq[1:]:
            k = _pptx_rpr_key(rPr)
            if k == cur_key:
                cur_chars.append(ch)
            else:
                runs.append((''.join(cur_chars), cur_rPr))
                cur_key, cur_rPr, cur_chars = k, rPr, [ch]
        runs.append((''.join(cur_chars), cur_rPr))

    # Remove existing <a:r> elements; preserve <a:pPr> and other non-run children
    for r in list(ap_elem.findall(_A + 'r')):
        ap_elem.remove(r)

    # Find insertion point: after <a:pPr> if present
    insert_idx = 0
    for idx, child in enumerate(ap_elem):
        if child.tag == _A + 'pPr':
            insert_idx = idx + 1
            break

    for offset, (text, rPr) in enumerate(runs):
        ap_elem.insert(insert_idx + offset, _pptx_make_run(text, rPr))

    return True


# ── Main writeback function ────────────────────────────────────────────────────

def writeback_pptx(new_md_path, base_md_path, pptx_path):
    """Apply text changes (new_md vs base_md) to source.pptx using shapeId/paraIndex anchors.
    Modifies pptx_path in-place. Returns {ok, changed, skipped, errors}."""
    from pptx import Presentation

    with open(new_md_path,  encoding='utf-8') as f:
        new_md  = f.read()
    with open(base_md_path, encoding='utf-8') as f:
        base_md = f.read()

    new_paras  = parse_pptx_md_anchors(new_md)
    base_paras = parse_pptx_md_anchors(base_md)

    # Find changed paragraphs
    changed = {}   # {composite_key: (old_plain, new_plain)}
    for key, new_text in new_paras.items():
        old_text = base_paras.get(key)
        if old_text is None:
            continue
        if new_text == old_text:
            continue
        old_plain = _strip_md_markup(old_text)
        new_plain = _strip_md_markup(new_text)
        if old_plain == new_plain:
            continue
        changed[key] = (old_plain, new_plain)

    if not changed:
        return {'ok': True, 'changed': 0, 'skipped': 0, 'errors': []}

    prs = Presentation(pptx_path)

    # Build lookup: {(slide_num, shape_id): shape} for all text-frame shapes
    shape_index = {}
    for sn, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_index[(sn, shape.shape_id)] = shape

    n_changed = 0
    n_skipped = 0
    errors    = []

    for key, (old_plain, new_plain) in changed.items():
        # key = "slideN:shapeId/paraIdx"
        try:
            colon_pos = key.index(':')
            slide_num = int(key[:colon_pos])
            rest      = key[colon_pos + 1:]   # "shapeId/paraIdx"
            slash_pos = rest.index('/')
            shape_id  = int(rest[:slash_pos])
            para_idx  = int(rest[slash_pos + 1:])
        except (ValueError, IndexError) as e:
            errors.append({'key': key, 'error': 'key parse error: ' + str(e)})
            continue

        shape = shape_index.get((slide_num, shape_id))
        if shape is None:
            n_skipped += 1
            errors.append({
                'key': key,
                'note': 'shape not found (slide %d, shape_id %d)' % (slide_num, shape_id),
            })
            continue

        tf = shape.text_frame
        if para_idx >= len(tf.paragraphs):
            n_skipped += 1
            errors.append({
                'key': key,
                'note': 'paraIndex %d out of range (%d paras)' % (para_idx, len(tf.paragraphs)),
            })
            continue

        para   = tf.paragraphs[para_idx]
        ap_elem = para._p   # direct lxml <a:p> element

        try:
            did_change = _rewrite_pptx_para(ap_elem, new_plain)
            if did_change:
                n_changed += 1
            else:
                n_skipped += 1
        except Exception as e:
            errors.append({'key': key, 'error': str(e)})

    if n_changed > 0:
        prs.save(pptx_path)

    return {
        'ok':      True,
        'changed':  n_changed,
        'skipped':  n_skipped,
        'errors':   errors,
    }


# ── P2-PPTX (DOC-P2P-01): in-place text edits keyed directly by shape/para anchor ──
# The realistic-preview in-place editor edits the pptx-renderer DOM's <div data-para-idx>
# directly, so it already knows {compositeKey: newPlainText} — no md diff/alignment needed.
# Reuses _rewrite_pptx_para (only changed chars move; every other run's format is preserved).
# Composite key formats (slide is 1-based, matching enumerate(prs.slides, start=1)):
#   text frame : "<slide>:<shapeId>/<paraIdx>"            e.g. "1:5/0"
#   table cell : "<slide>:<shapeId>/<row>/<col>/<paraIdx>" e.g. "1:7/2/1/0"
# (Table cells need row/col because a table's cells share one shapeId and each cell's
#  paragraph index restarts at 0 — shapeId/paraIdx alone would collide across cells.)

def _parse_edit_key(key):
    """Parse a composite edit key → (slide_num, shape_id, row, col, para_idx).
    row/col are None for text-frame (non-table) anchors. Raises ValueError on malformed key."""
    colon = key.index(':')
    slide_num = int(key[:colon])
    rest      = key[colon + 1:]
    parts     = rest.split('/')
    if len(parts) == 2:                 # shapeId/paraIdx  (text frame)
        return slide_num, int(parts[0]), None, None, int(parts[1])
    if len(parts) == 4:                 # shapeId/row/col/paraIdx  (table cell)
        return slide_num, int(parts[0]), int(parts[1]), int(parts[2]), int(parts[3])
    raise ValueError('unexpected key shape: ' + key)


def writeback_edits(pptx_path, edits):
    """Apply {compositeKey: new_plain_text} directly to source.pptx via run-level rewrite.
    Only paragraphs whose text actually differs are touched. Returns
    {ok, changed, skipped, errors}."""
    from pptx import Presentation

    edits = {str(k): (v if v is not None else '') for k, v in (edits or {}).items()}
    if not edits:
        return {'ok': True, 'changed': 0, 'skipped': 0, 'errors': []}

    prs = Presentation(pptx_path)

    # Index every shape by (slide_num, shape_id) — one id is unique within a slide.
    shape_index = {}
    for sn, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            shape_index[(sn, shape.shape_id)] = shape

    changed = 0
    skipped = 0
    errors  = []

    for key, new_text in edits.items():
        try:
            slide_num, shape_id, row, col, para_idx = _parse_edit_key(key)
        except (ValueError, IndexError) as e:
            errors.append({'key': key, 'error': 'key parse error: ' + str(e)})
            continue

        shape = shape_index.get((slide_num, shape_id))
        if shape is None:
            skipped += 1
            errors.append({'key': key, 'note': 'shape not found (slide %d, shape_id %d)' % (slide_num, shape_id)})
            continue

        # Resolve the target text_frame (shape body, or a table cell).
        try:
            if row is not None:
                if not getattr(shape, 'has_table', False):
                    skipped += 1
                    errors.append({'key': key, 'note': 'shape %d is not a table' % shape_id})
                    continue
                table = shape.table
                if row >= len(table.rows) or col >= len(table.columns):
                    skipped += 1
                    errors.append({'key': key, 'note': 'cell (%d,%d) out of range' % (row, col)})
                    continue
                tf = table.cell(row, col).text_frame
            else:
                if not getattr(shape, 'has_text_frame', False):
                    skipped += 1
                    errors.append({'key': key, 'note': 'shape %d has no text frame' % shape_id})
                    continue
                tf = shape.text_frame
        except Exception as e:
            errors.append({'key': key, 'error': 'resolve text frame: ' + str(e)})
            continue

        if para_idx >= len(tf.paragraphs):
            skipped += 1
            errors.append({'key': key, 'note': 'paraIndex %d out of range (%d paras)' % (para_idx, len(tf.paragraphs))})
            continue

        ap_elem = tf.paragraphs[para_idx]._p
        try:
            if _rewrite_pptx_para(ap_elem, new_text):
                changed += 1
            else:
                skipped += 1
        except Exception as e:
            errors.append({'key': key, 'error': str(e)})

    if changed > 0:
        prs.save(pptx_path)

    return {'ok': True, 'changed': changed, 'skipped': skipped, 'errors': errors}


if __name__ == '__main__':
    ap = argparse.ArgumentParser(description='K4-02: pptx run-level writeback')
    ap.add_argument('--pptx', required=True, help='path to source.pptx (modified in-place)')
    ap.add_argument('--new',  help='path to new doc.md (md-diff mode; with --base)')
    ap.add_argument('--base', help='path to base doc.md (md-diff mode; with --new)')
    ap.add_argument('--edits', help='path to JSON {compositeKey: newPlainText} (in-place mode, DOC-P2P-01)')
    args = ap.parse_args()

    if args.edits:
        with open(args.edits, encoding='utf-8') as f:
            edits = json.load(f)
        result = writeback_edits(args.pptx, edits)
    elif args.new and args.base:
        result = writeback_pptx(args.new, args.base, args.pptx)
    else:
        result = {'ok': False, 'error': 'must pass either --edits, or both --new and --base'}

    print(json.dumps(result, ensure_ascii=False))
    sys.exit(0 if result.get('ok') else 1)
